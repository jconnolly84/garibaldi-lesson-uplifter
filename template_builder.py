# template_builder.py

from pptx import Presentation
from pptx.util import Inches, Pt
import tempfile
import os
import requests

from app import fetch_best_image, fetch_youtube_video

def insert_images_into_template(uplifted_text):
    prs = Presentation("Lessonary_Template_Structure.pptx")
    blank_layout = prs.slide_layouts[6]  # blank slide layout

    # Cleanup existing slides
    for i in range(len(prs.slides) - 1, -1, -1):
        r_id = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[i]

    slide_blocks = uplifted_text.strip().split("--- Slide")

    for block in slide_blocks:
        if not block.strip():
            continue

        lines = block.strip().split("\n")
        slide_title_line = lines[0].strip(": ")
        body_lines = lines[1:]
        content = "\n".join(body_lines).strip()

        # Add slide
        slide = prs.slides.add_slide(blank_layout)

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = slide_title_line
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.name = "Century Gothic"

        # Add main content
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.5), Inches(4.5))
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.text = content
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.name = "Century Gothic"

        # Add image
        image_url = fetch_best_image(content)
        if image_url:
            try:
                image_data = requests.get(image_url).content
                tmp_img = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
                tmp_img.write(image_data)
                tmp_img.close()
                slide.shapes.add_picture(tmp_img.name, Inches(6.2), Inches(1.3), width=Inches(3.5))
                os.unlink(tmp_img.name)
            except:
                pass

        # Add YouTube link
        video_url = fetch_youtube_video(content)
        if video_url:
            video_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.5))
            video_tf = video_box.text_frame
            p = video_tf.paragraphs[0]
            run = p.add_run()
            run.text = "📺 Click to watch: YouTube Video"
            run.hyperlink.address = video_url
            run.font.size = Pt(14)
            run.font.name = "Century Gothic"

    return prs
