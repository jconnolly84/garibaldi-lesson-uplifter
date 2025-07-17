# Lessonary Uplifter (Streamlit GUI Version with Template, Strategy, and Media Support + AI Boost Option)

import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from openai import OpenAI
import io
from PIL import Image
import requests
import tempfile
import os
from datetime import datetime

# --- Config ---
client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])

# --- Image Fetching Functions ---
def fetch_pixabay_image(query):
    api_key = st.secrets["PIXABAY_API_KEY"]
    url = f"https://pixabay.com/api/?key={api_key}&q={query}&image_type=photo&safesearch=true&per_page=3"
    response = requests.get(url)
    if response.status_code == 200:
        data = response.json()
        if data["hits"]:
            return data["hits"][0]["webformatURL"]
    return None

def fetch_pexels_image(query):
    api_key = st.secrets["PEXELS_API_KEY"]
    headers = {"Authorization": api_key}
    url = f"https://api.pexels.com/v1/search?query={query}&per_page=3"
    response = requests.get(url, headers=headers)
    if response.status_code == 200:
        data = response.json()
        if data["photos"]:
            return data["photos"][0]["src"]["medium"]
    return None

def fetch_best_image(query):
    query = query.replace("\n", " ").strip()
    image_url = fetch_pixabay_image(query)
    if not image_url:
        image_url = fetch_pexels_image(query)
    return image_url

# --- YouTube Fetching Function ---
def fetch_youtube_video(query):
    api_key = st.secrets["YOUTUBE_API_KEY"]
    search_url = f"https://www.googleapis.com/youtube/v3/search?part=snippet&q={query}&type=video&maxResults=1&key={api_key}"
    response = requests.get(search_url)
    if response.status_code == 200:
        items = response.json().get("items", [])
        if items:
            video_id = items[0]["id"]["videoId"]
            return f"https://www.youtube.com/watch?v={video_id}"
    return None

# --- Prompt Builder ---
def build_prompt(slide_text, enrichment_level="Base"):
    extra = ""
    if enrichment_level == "Enhanced":
        extra = "Add real-world links, extend explanations, and enrich slides with more detail."
    elif enrichment_level == "Max":
        extra = "Add interactive tasks, teacher narration prompts, cross-curricular links, and advanced vocabulary."

    prompt = f"""
You are an expert teacher and lesson designer at a secondary school. A teacher has uploaded a PowerPoint lesson.

Please analyse and rebuild the lesson using the following structure:

1. Ready to Learn / Entry
2. Connect & Recall
3. Explore / Impart Knowledge
4. Collaborate / Facilitate
5. Independent Practice (FIT)
6. Review & Improve
7. Homework

Your task:
- Reorder content into that structure
- Suggest new slides where needed (title + content)
- Improve clarity, challenge, and engagement
- Recommend relevant images or diagrams for each slide
- Recommend a supporting YouTube video only if it enhances learning
- Embed school-wide TLC strategies from TLC_Strategies.txt
- Include a Key Objective and Differentiated Outcomes slide
- Include a Vocabulary slide (max 6 terms)
- Include a What is the Connection slide with 4 image prompts
- End with a Homework task slide (relevant extension task, max 30 mins, bring into next lesson)
{extra}

Here is the raw content:

{slide_text}

Return the uplifted slide-by-slide version, labelled with headers like:
--- Slide 1: Ready to Learn ---
[Slide title, content, image suggestion, optional video link, strategy used]
"""
    return prompt

# --- GPT API Call ---
def call_chatgpt(prompt):
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.5
    )
    return response.choices[0].message.content

# --- Extract Text ---
def extract_text_from_pptx(file):
    prs = Presentation(file)
    content = []
    for i, slide in enumerate(prs.slides):
        slide_text = f"Slide {i+1}:\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
        content.append(slide_text.strip())
    return "\n\n".join(content)

# --- Slide Generator ---
def generate_pptx_from_text(uplifted_text):
    prs = Presentation()
    for block in uplifted_text.split("--- Slide"):
        if not block.strip():
            continue
        lines = block.strip().split("\n")
        if len(lines) < 2:
            continue
        title = lines[0].replace(":", "").strip()
        content = "\n".join(lines[1:]).strip()

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        title_tf = title_box.text_frame
        title_tf.text = title

        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(8.5), Inches(5))
        content_tf = content_box.text_frame
        content_tf.text = content

        image_url = fetch_best_image(content)
        if image_url:
            img_data = requests.get(image_url).content
            tmp_img = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
            tmp_img.write(img_data)
            tmp_img.close()
            slide.shapes.add_picture(tmp_img.name, Inches(6), Inches(1.5), width=Inches(3))
            os.unlink(tmp_img.name)

        video_url = fetch_youtube_video(content)
        if video_url:
            link_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(6), Inches(0.4))
            link_tf = link_box.text_frame
            p = link_tf.paragraphs[0]
            run = p.add_run()
            run.text = "ðŸ“º Watch Video"
            run.hyperlink.address = video_url
            run.font.size = Pt(12)

    return prs

# --- Streamlit App ---
st.set_page_config(page_title="Lessonary Uplifter", layout="centered")
st.image("LOGO_Lessonary.png", use_container_width=True)
st.title("Lessonary Uplifter")
st.write("ðŸ“š Upload a PowerPoint and get it automatically restructured using your school's T&L model.")

enrichment_level = st.selectbox("Optional AI Boost: Enrichment Level", ["Base", "Enhanced", "Max"], index=0)

uploaded_file = st.file_uploader("Upload a .pptx file", type="pptx")

if uploaded_file is not None:
    uploaded_filename = uploaded_file.name.replace(".pptx", "")
    upload_date = datetime.now().strftime("%Y%m%d")
    output_filename = f"{uploaded_filename}_uplifted_{upload_date}.pptx"

    with st.spinner("Extracting slide text and analysing..."):
        slide_text = extract_text_from_pptx(uploaded_file)
        prompt = build_prompt(slide_text, enrichment_level=enrichment_level)
        uplifted_lesson = call_chatgpt(prompt)

    st.success("âœ… Lesson uplift complete!")
    st.subheader("ðŸ” Uplifted Lesson Structure")
    st.text_area("Slide-by-slide output:", uplifted_lesson, height=600)

    st.download_button(
        label="ðŸ“¥ Download as text file",
        data=uplifted_lesson,
        file_name="Lessonary_Uplifted_Lesson.txt",
        mime="text/plain"
    )

    if st.button("ðŸ“¤ Download as PPTX with AI Images & Videos"):
        pptx_output = generate_pptx_from_text(uplifted_lesson)
        tmp_pptx = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        pptx_output.save(tmp_pptx.name)
        tmp_pptx.close()
        with open(tmp_pptx.name, "rb") as f:
            st.download_button(
                label="ðŸ“¥ Download PowerPoint File",
                data=f,
                file_name=output_filename,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        os.unlink(tmp_pptx.name)
